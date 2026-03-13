from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import io

def create_pitch_deck(business_plan):
    """
    Takes a JSON business plan and generates a 10-slide Pitch Deck (PPTX format).
    Returns the binary content of the PPTX file.
    """
    prs = Presentation()
    
    # Extract nested dictionaries safely
    idea = business_plan.get('idea', {})
    location = business_plan.get('location', {})
    pricing = business_plan.get('pricing', {})
    marketing = business_plan.get('marketing', {})
    financials = business_plan.get('financials', {})
    growth = business_plan.get('growth', {})

    
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0]) # Title Slide layout
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = idea.get('name', 'Startup Pitch Deck')
    subtitle.text = idea.get('description', 'A New Venture')
    
    # Slide 2: Problem
    slide = prs.slides.add_slide(prs.slide_layouts[1]) # Title and Content
    slide.shapes.title.text = "The Problem"
    content = slide.placeholders[1]
    content.text = f"Current solutions in the market for {idea.get('name', 'this industry')} are lacking or too expensive."

    # Slide 3: Solution (Product Overview)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "The Solution / Product Overview"
    content = slide.placeholders[1]
    content.text = idea.get('description', 'Our unique product/service offering.')

    # Slide 4: Market Opportunity
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Market Opportunity"
    content = slide.placeholders[1]
    content.text = f"Targeting {location.get('areaType', 'selected region')} with an estimated shop size of {location.get('shopSize', 'standard')}."

    # Slide 5: Business Model
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Business Model"
    content = slide.placeholders[1]
    content.text = "Revenue streams, pricing strategy, and cost structure."
    
    p_frame = content.text_frame
    if pricing.get('suggestedPrice'):
        p = p_frame.add_paragraph()
        p.text = f"Suggested Price: {pricing.get('suggestedPrice')} with a Margin of {pricing.get('profitMargin', 'standard')}."

    # Slide 6: Competitor Landscape
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Competitor Landscape"
    content = slide.placeholders[1]
    content.text = "We differentiate ourselves through unique value propositions and superior service."

    # Slide 7: Go-To-Market Strategy
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Go-To-Market Strategy"
    content = slide.placeholders[1]
    
    if marketing.get('launchPlan'):
        channels_text = ", ".join(marketing.get('launchPlan', []))
        content.text = f"Key Marketing Channels: {channels_text}"
    else:
        content.text = "Social media, local partnerships, and targeted advertising."

    # Slide 8: Financial Projection
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Financial Projections"
    content = slide.placeholders[1]
    
    content.text = f"Estimated Initial Cost: {idea.get('investmentRange', 'N/A')}\n"
    content.text += f"Projected Monthly Revenue: {idea.get('expectedRevenue', 'N/A')}\n"
    content.text += f"Expected Profit Margin: {idea.get('profitMargin', 'N/A')}\n"
    content.text += f"Break-even Timeline: {idea.get('breakEvenTime', 'N/A')}"

    # Slide 9: Funding Requirement
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Funding Requirement"
    content = slide.placeholders[1]
    budget_req = idea.get('investmentRange', 'the required capital')
    content.text = f"Seeking {budget_req} to launch operations and reach profitability."
    
    # Slide 10: Vision and Expansion Plan
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Vision and Expansion Plan"
    content = slide.placeholders[1]
    
    if growth.get('sixMonths'):
        content.text = f"6-Month Goal: {growth.get('sixMonths', '')}"
    else:
        content.text = "Expand footprint, increase product line, and capture regional market share."

    # Save to memory stream
    pptx_stream = io.BytesIO()
    prs.save(pptx_stream)
    pptx_stream.seek(0)
    
    return pptx_stream.read()

def create_bank_loan_pdf(business_plan):
    """
    Generates a standardized Bank Loan Application format (e.g. MSME / MUDRA) as a PDF using ReportLab
    Returns the binary content of the PDF file.
    """
    from reportlab.lib.pagesizes import letter
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib import colors

    idea = business_plan.get('idea', {})
    location = business_plan.get('location', {})
    pricing = business_plan.get('pricing', {})

    pdf_stream = io.BytesIO()
    doc = SimpleDocTemplate(pdf_stream, pagesize=letter)
    
    styles = getSampleStyleSheet()
    title_style = styles['Heading1']
    title_style.alignment = 1 # Center
    
    normal_style = styles['Normal']
    
    story = []
    
    # Header
    story.append(Paragraph("MUDRA / MSME LOAN APPLICATION FORMAT", title_style))
    story.append(Spacer(1, 20))
    
    story.append(Paragraph(f"<b>Business Name:</b> {idea.get('name', 'Pending Registration')}", normal_style))
    story.append(Paragraph(f"<b>Proposed Location Area:</b> {location.get('areaType', 'TBD')} ({location.get('shopSize', 'N/A')})", normal_style))
    story.append(Paragraph(f"<b>Business Activity:</b> {idea.get('description', 'N/A')}", normal_style))
    story.append(Spacer(1, 15))
    
    # Financials Table
    data = [
        ['Financial Indicator', 'Amount / Value'],
        ['Total Initial Investment Required', f"{idea.get('investmentRange', '0')}"],
        ['Projected Monthly Revenue', f"{idea.get('expectedRevenue', '0')}"],
        ['Estimated Cost Price per unit', f"{pricing.get('costPrice', '0')}"],
        ['Profit Margin', f"{idea.get('profitMargin', '0')}"]
    ]
    
    t = Table(data, colWidths=[250, 200])
    t.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,0), colors.HexColor("#2a1b54")),
        ('TEXTCOLOR', (0,0), (-1,0), colors.whitesmoke),
        ('ALIGN', (0,0), (-1,-1), 'LEFT'),
        ('FONTNAME', (0,0), (-1,0), 'Helvetica-Bold'),
        ('BOTTOMPADDING', (0,0), (-1,0), 12),
        ('BACKGROUND', (0,1), (-1,-1), colors.HexColor("#f8f9fa")),
        ('GRID', (0,0), (-1,-1), 1, colors.lightgrey)
    ]))
    
    story.append(Paragraph("<b>Financial Summary</b>", styles['Heading2']))
    story.append(Spacer(1, 10))
    story.append(t)
    story.append(Spacer(1, 20))

    story.append(Paragraph("<b>Declaration:</b>", styles['Heading3']))
    story.append(Paragraph("I/We hereby declare that all information provided in this preliminary project report is true and correct to the best of my/our knowledge.", normal_style))
    story.append(Spacer(1, 40))
    
    story.append(Paragraph("Signature: ___________________________", normal_style))
    story.append(Spacer(1, 10))
    story.append(Paragraph("Date: ___________________________", normal_style))

    doc.build(story)
    
    pdf_stream.seek(0)
    return pdf_stream.read()
